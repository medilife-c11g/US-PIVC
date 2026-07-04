#!/usr/bin/env python3
"""Build the eFAST trauma keynote deck (16:9 PPTX).

Usage: python3 build_deck.py
Output: eFAST_trauma_keynote.pptx (same folder)
"""

import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- palette
NAVY = RGBColor(0x0B, 0x27, 0x40)      # primary dark
TEAL = RGBColor(0x14, 0x91, 0x9B)      # accent
TEAL_DARK = RGBColor(0x0E, 0x6B, 0x72)
AMBER = RGBColor(0xF2, 0xA5, 0x41)     # highlight
RED = RGBColor(0xC9, 0x44, 0x4D)       # warnings / pitfalls
SLATE = RGBColor(0x35, 0x42, 0x4E)     # body text
GRAY = RGBColor(0x6E, 0x7B, 0x87)      # muted
LIGHT = RGBColor(0xF4, 0xF7, 0xF9)     # light panel fill
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Segoe UI"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

slide_counter = {"n": 0}


# ---------------------------------------------------------------- helpers
def _set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _dashed_outline(shape, color=GRAY, width=Pt(1.25)):
    shape.fill.background()
    shape.line.color.rgb = color
    shape.line.width = width
    ln = shape.line._get_or_add_ln()
    dash = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
    ln.append(dash)


def _text(shape, runs, size=18, color=SLATE, bold=False, align=PP_ALIGN.LEFT,
          anchor=MSO_ANCHOR.TOP, line_spacing=1.0, space_after=0):
    """runs: str, or list of (text, dict-of-overrides) tuples for one paragraph."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    if space_after:
        p.space_after = Pt(space_after)
    if isinstance(runs, str):
        runs = [(runs, {})]
    for txt, ov in runs:
        r = p.add_run()
        r.text = txt
        f = r.font
        f.name = ov.get("font", FONT)
        f.size = Pt(ov.get("size", size))
        f.bold = ov.get("bold", bold)
        f.italic = ov.get("italic", False)
        f.color.rgb = ov.get("color", color)
    return tf


def _box(slide, x, y, w, h):
    return slide.shapes.add_textbox(x, y, w, h)


def add_footer(slide, section=""):
    slide_counter["n"] += 1
    n = slide_counter["n"]
    if n == 1:
        return
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), SLIDE_H - Inches(0.05),
                                 SLIDE_W, Inches(0.05))
    _set_fill(bar, TEAL)
    tb = _box(slide, Inches(0.35), SLIDE_H - Inches(0.42), Inches(9), Inches(0.32))
    _text(tb, f"eFAST in Trauma — Keynote{('  ·  ' + section) if section else ''}",
          size=10, color=GRAY)
    num = _box(slide, SLIDE_W - Inches(1.0), SLIDE_H - Inches(0.42), Inches(0.7), Inches(0.32))
    _text(num, str(n), size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def content_slide(title, kicker="", section=""):
    """Standard white content slide: teal kicker, navy title, accent rule."""
    slide = prs.slides.add_slide(BLANK)
    if kicker:
        kb = _box(slide, Inches(0.55), Inches(0.32), Inches(11), Inches(0.35))
        _text(kb, kicker.upper(), size=13, color=TEAL, bold=True)
    tb = _box(slide, Inches(0.5), Inches(0.62), Inches(12.3), Inches(0.95))
    _text(tb, title, size=32, color=NAVY, bold=True)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.58), Inches(1.52),
                                  Inches(1.6), Inches(0.06))
    _set_fill(rule, AMBER)
    add_footer(slide, section)
    return slide


def add_bullets(slide, items, x=Inches(0.6), y=Inches(1.85), w=Inches(12.1),
                h=Inches(5.2), size=19, gap=10):
    """items: list of (level, text) or (level, [(run, overrides), ...])."""
    tb = _box(slide, x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        level, content = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
        runs = [(content, {})] if isinstance(content, str) else content
        marker = "▪  " if level == 0 else "–  "
        mr = p.add_run()
        mr.text = marker
        mr.font.name = FONT
        mr.font.size = Pt(size if level == 0 else size - 2)
        mr.font.color.rgb = TEAL if level == 0 else GRAY
        mr.font.bold = level == 0
        for txt, ov in runs:
            r = p.add_run()
            r.text = txt
            f = r.font
            f.name = FONT
            f.size = Pt(ov.get("size", size if level == 0 else size - 2))
            f.bold = ov.get("bold", False)
            f.italic = ov.get("italic", False)
            f.color.rgb = ov.get("color", SLATE if level == 0 else GRAY)
    return tb


def image_placeholder(slide, x, y, w, h, label):
    ph = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    _dashed_outline(ph)
    _text(ph, [("INSERT IMAGE / CLIP\n", {"bold": True, "size": 13, "color": GRAY}),
               (label, {"size": 12, "color": GRAY, "italic": True})],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return ph


def panel(slide, x, y, w, h, fill=LIGHT, accent=None):
    p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    _set_fill(p, fill)
    if accent:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.07), h)
        _set_fill(bar, accent)
    return p


def panel_text(slide, x, y, w, h, header, body, fill=LIGHT, accent=TEAL,
               header_color=None, body_size=15):
    panel(slide, x, y, w, h, fill=fill, accent=accent)
    tb = _box(slide, x + Inches(0.22), y + Inches(0.12), w - Inches(0.4), h - Inches(0.24))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = header
    r.font.name = FONT
    r.font.size = Pt(body_size + 2)
    r.font.bold = True
    r.font.color.rgb = header_color or accent
    for line in body if isinstance(body, list) else [body]:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        p2.line_spacing = 1.05
        r2 = p2.add_run()
        r2.text = line
        r2.font.name = FONT
        r2.font.size = Pt(body_size)
        r2.font.color.rgb = SLATE
    return tb


def divider(number, title, subtitle, notes):
    slide = prs.slides.add_slide(BLANK)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _set_fill(bg, NAVY)
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.9), SLIDE_W, Inches(0.12))
    _set_fill(stripe, TEAL)
    nb = _box(slide, Inches(0.9), Inches(1.7), Inches(3), Inches(1.6))
    _text(nb, number, size=110, color=TEAL, bold=True)
    tb = _box(slide, Inches(0.95), Inches(3.45), Inches(11.9), Inches(1.1))
    _text(tb, title, size=44, color=WHITE, bold=True)
    sb = _box(slide, Inches(0.98), Inches(4.6), Inches(11.2), Inches(0.9))
    _text(sb, subtitle, size=20, color=RGBColor(0xB8, 0xC9, 0xD9))
    add_footer(slide)
    add_notes(slide, notes)
    return slide


# ================================================================ SLIDE 1 — title
slide = prs.slides.add_slide(BLANK)
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
_set_fill(bg, NAVY)
side = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.28), SLIDE_H)
_set_fill(side, TEAL)
kb = _box(slide, Inches(1.0), Inches(1.35), Inches(11), Inches(0.5))
_text(kb, "KEYNOTE", size=16, color=AMBER, bold=True)
tb = _box(slide, Inches(0.95), Inches(1.9), Inches(11.8), Inches(2.2))
_text(tb, "eFAST in Trauma:\nFrom the Resus Bay to the Roadside", size=40, color=WHITE, bold=True)
sb = _box(slide, Inches(1.0), Inches(4.05), Inches(11.4), Inches(0.7))
_text(sb, "Extended Focused Assessment with Sonography for Trauma —"
          " applications in everyday, rural, and mass-casualty care",
      size=21, color=RGBColor(0xB8, 0xC9, 0xD9))
ab = _box(slide, Inches(1.0), Inches(5.6), Inches(11), Inches(1.2))
_text(ab, [("Chia-Ching Chen, MD\n", {"size": 22, "bold": True, "color": WHITE}),
           ("Emergency Medicine, Changhua Show Chwan Memorial Hospital\n",
            {"size": 15, "color": RGBColor(0xB8, 0xC9, 0xD9)}),
           ("[DATE / VENUE]", {"size": 15, "color": TEAL, "italic": True})])
add_footer(slide)
add_notes(slide, "Welcome. One-hour keynote: ~45 min talk, 10 min Q&A, 5 min buffer. "
                 "Framing: eFAST is not just a trauma-bay ritual — it is a decision "
                 "instrument that matters MOST where resources are thinnest: rural "
                 "hospitals and mass-casualty scenes. Update date/venue before presenting.")

# ================================================================ SLIDE 2 — disclosures
slide = content_slide("Disclosures")
add_bullets(slide, [
    (0, "[No conflicts of interest to declare — EDIT AS APPLICABLE]"),
    (0, "[Funding / device relationships, if any]"),
    (0, "Images and clips in this talk are used for education. Patient-identifiable "
        "material has been removed or consented."),
])
add_notes(slide, "Keep brief — 15 seconds. Edit to reflect your actual disclosures "
                 "and your institution's/conference's requirements.")

# ================================================================ SLIDE 3 — roadmap
slide = content_slide("Where we are going", kicker="Roadmap")
items = [
    ("1", "Foundations", "What eFAST asks, the views, the evidence, the pitfalls"),
    ("2", "Decision-making", "Where eFAST sits in the trauma algorithm"),
    ("3", "The rural reality", "Triage-to-transfer when there is no CT down the hall"),
    ("4", "Mass casualty", "Sonographic triage when needs exceed resources"),
    ("5", "Looking ahead", "Training, tele-ultrasound, AI — and 5 take-homes"),
]
x0, y0, w, h = Inches(0.6), Inches(1.95), Inches(12.1), Inches(0.92)
for i, (num, head, sub) in enumerate(items):
    y = y0 + i * Inches(1.02)
    panel(slide, x0, y, w, h, fill=LIGHT, accent=TEAL if i % 2 == 0 else NAVY)
    nb = _box(slide, x0 + Inches(0.25), y + Inches(0.12), Inches(0.7), Inches(0.7))
    _text(nb, num, size=30, color=AMBER, bold=True)
    tb = _box(slide, x0 + Inches(1.1), y + Inches(0.08), Inches(3.6), Inches(0.75))
    _text(tb, head, size=20, color=NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    sb = _box(slide, x0 + Inches(4.9), y + Inches(0.08), Inches(7.0), Inches(0.78))
    _text(sb, sub, size=15, color=GRAY, anchor=MSO_ANCHOR.MIDDLE)
add_notes(slide, "~1 min. Signpost the hour. Emphasize the arc: we start in the "
                 "resus bay everyone knows, then take the same probe to places "
                 "with no CT scanner and no backup.")

# ================================================================ SLIDE 4 — opening case
slide = content_slide("A Tuesday night, 90 minutes from CT",
                      kicker="Opening case", section="Case")
add_bullets(slide, [
    (0, [("34-year-old motorcyclist", {"bold": True}),
         (" vs. truck, brought to a small rural hospital", {})]),
    (1, "BP 86/54, HR 128, RR 26, GCS 14; abdomen tender, seatbelt-type abrasion"),
    (1, "No CT scanner tonight. No surgeon in-house. Helicopter: 40 min out."),
    (0, [("The only imaging you have is the ultrasound machine beside you.", {"bold": True})]),
    (0, [("Question for the room: ", {"bold": True, "color": TEAL}),
         ("what single examination changes what you do in the next 10 minutes?", {})]),
], y=Inches(1.9), h=Inches(3.4))
image_placeholder(slide, Inches(8.9), Inches(4.55), Inches(3.8), Inches(2.35),
                  "Optional: photo of a rural ED / transport map")
add_notes(slide, "~2 min. Hook the audience with the case — pause on the question. "
                 "We return to this patient at the end of Part 2. The point: the "
                 "value of eFAST is inversely proportional to the resources around you.")

# ================================================================ PART 1
divider("01", "Foundations of eFAST",
        "Four questions, ten minutes of anatomy, and an honest look at the evidence",
        "~30 sec transition. Even experts benefit from re-anchoring on what the "
        "exam can and cannot answer — the rural and MCI material later depends on it.")

# ---- 5: history
slide = content_slide("From FAST to eFAST — a 30-year evolution",
                      kicker="Part 1 · Foundations", section="Foundations")
steps = [
    ("1970s–80s", "European surgeons use bedside ultrasound for blunt abdominal "
                  "trauma, largely replacing peritoneal lavage"),
    ("1990s", "North American adoption; the term “FAST” is standardized through "
              "international consensus (Rozycki, Scalea, and colleagues)"),
    ("2004", "“Extended” FAST described: thoracic windows added for pneumothorax "
             "and hemothorax (Kirkpatrick et al.)"),
    ("Today", "An adjunct to the ATLS® primary survey — performed during, "
              "not after, resuscitation"),
]
y0 = Inches(1.95)
for i, (era, txt) in enumerate(steps):
    y = y0 + i * Inches(1.18)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.75), y + Inches(0.18),
                                 Inches(0.28), Inches(0.28))
    _set_fill(dot, TEAL if i < 3 else AMBER)
    if i < 3:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.86), y + Inches(0.5),
                                      Inches(0.045), Inches(0.85))
        _set_fill(line, RGBColor(0xC8, 0xD4, 0xDE))
    eb = _box(slide, Inches(1.3), y, Inches(1.9), Inches(0.6))
    _text(eb, era, size=17, color=NAVY, bold=True)
    tb = _box(slide, Inches(3.3), y - Inches(0.03), Inches(9.4), Inches(1.1))
    _text(tb, txt, size=16, color=SLATE)
add_notes(slide, "~2 min. Key historical beats only — verify the exact citation "
                 "details against the reference slide before presenting. The message: "
                 "each expansion of the exam followed a clinical need, not a "
                 "technology push. eFAST added the chest because supine CXR misses "
                 "so many pneumothoraces.")

# ---- 6: four questions
slide = content_slide("eFAST asks four binary questions", kicker="Part 1 · Foundations",
                      section="Foundations")
qs = [
    ("Free intraperitoneal fluid?", "Hemoperitoneum until proven otherwise in trauma", TEAL),
    ("Pericardial effusion?", "± sonographic signs of tamponade", NAVY),
    ("Pneumothorax?", "Absent lung sliding, lung point", TEAL),
    ("Hemothorax?", "Fluid above the diaphragm", NAVY),
]
for i, (q, sub, c) in enumerate(qs):
    x = Inches(0.6) + (i % 2) * Inches(6.25)
    y = Inches(2.0) + (i // 2) * Inches(1.85)
    panel_text(slide, x, y, Inches(5.95), Inches(1.6), q, sub, accent=c, body_size=15)
tb = _box(slide, Inches(0.62), Inches(5.95), Inches(12.0), Inches(1.0))
_text(tb, [("Not ", {"size": 18, "color": SLATE}),
           ("“is there an injury?”", {"size": 18, "color": SLATE, "italic": True}),
           (" — but ", {"size": 18, "color": SLATE}),
           ("“does this patient need an intervention, an operating room, "
            "or a transfer — now?”", {"size": 18, "color": TEAL, "bold": True})])
add_notes(slide, "~1.5 min. Hammer the framing: eFAST is a disposition tool, not a "
                 "diagnostic survey. It is deliberately insensitive to solid-organ "
                 "grade, hollow viscus, and retroperitoneum — that honesty is what "
                 "makes it safe to use.")

# ---- 7: the views
slide = content_slide("The exam: six windows, three to five minutes",
                      kicker="Part 1 · Foundations", section="Foundations")
views = [
    ("1 · Subxiphoid cardiac", "Pericardial fluid; RV collapse if tamponade"),
    ("2 · RUQ (Morison's)", "Hepatorenal recess, caudal liver tip, right lung base"),
    ("3 · LUQ (splenorenal)", "Perisplenic + subphrenic space, left lung base"),
    ("4 · Suprapubic", "Rectovesical space / pouch of Douglas"),
    ("5–6 · Anterior chest ×2", "Lung sliding at the most anterior point, supine"),
]
for i, (head, sub) in enumerate(views):
    y = Inches(1.95) + i * Inches(0.98)
    panel_text(slide, Inches(0.6), y, Inches(7.0), Inches(0.86), head, sub,
               accent=TEAL, body_size=13)
image_placeholder(slide, Inches(8.0), Inches(1.95), Inches(4.7), Inches(3.6),
                  "Torso diagram with probe positions (your preferred figure)")
pb = panel(slide, Inches(8.0), Inches(5.75), Inches(4.7), Inches(1.05), fill=LIGHT, accent=AMBER)
tb = _box(slide, Inches(8.25), Inches(5.87), Inches(4.3), Inches(0.85))
_text(tb, [("Probe: ", {"bold": True, "size": 14, "color": NAVY}),
           ("curvilinear or phased 2–5 MHz for all views; "
            "linear probe optional for pleura", {"size": 14, "color": SLATE})])
add_notes(slide, "~2 min. Sequence is flexible — in an unstable patient many start "
                 "at the heart. Stress that a complete exam is fast; incomplete "
                 "exams are where misses happen.")

# ---- 8: RUQ
slide = content_slide("RUQ view — the workhorse", kicker="Part 1 · Foundations",
                      section="Foundations")
add_bullets(slide, [
    (0, [("Morison's pouch", {"bold": True}),
         (" (hepatorenal recess): classic single most sensitive window for free "
          "fluid in the supine patient", {})]),
    (0, [("Sweep to the caudal liver tip", {"bold": True}),
         (" — often the first place small volumes collect; a “Morison's-only” "
          "glance under-performs", {})]),
    (0, "Slide cephalad above the diaphragm → right hemothorax check in the same window"),
    (0, [("Pitfall: ", {"bold": True, "color": RED}),
         ("perinephric fat and gallbladder can mimic or hide free fluid", {})]),
], w=Inches(7.4))
image_placeholder(slide, Inches(8.35), Inches(1.95), Inches(4.35), Inches(2.35),
                  "Normal RUQ clip")
image_placeholder(slide, Inches(8.35), Inches(4.45), Inches(4.35), Inches(2.35),
                  "Positive RUQ — anechoic stripe in Morison's pouch")
add_notes(slide, "~2 min. If you show only two clips in this talk, show a normal and "
                 "a positive RUQ side by side — the audience anchors on the contrast.")

# ---- 9: LUQ
slide = content_slide("LUQ view — more posterior than you think",
                      kicker="Part 1 · Foundations", section="Foundations")
add_bullets(slide, [
    (0, [("“Knuckles to the bed”", {"bold": True}),
         (" — probe more posterior and cephalad than the RUQ mirror image", {})]),
    (0, [("Check the subphrenic space first: ", {"bold": True}),
         ("fluid often collects between spleen and diaphragm before the "
          "splenorenal recess", {})]),
    (0, "Rib shadows dominate — small fanning movements, breath timing if feasible"),
    (0, "Same window, look above the diaphragm → left hemothorax"),
], w=Inches(7.4))
image_placeholder(slide, Inches(8.35), Inches(1.95), Inches(4.35), Inches(2.35),
                  "LUQ probe position photo")
image_placeholder(slide, Inches(8.35), Inches(4.45), Inches(4.35), Inches(2.35),
                  "Positive LUQ — perisplenic fluid")
add_notes(slide, "~1.5 min. The LUQ is the most commonly botched view; the "
                 "subphrenic-first habit is the single highest-yield teaching point.")

# ---- 10: pelvis
slide = content_slide("Pelvic view — the bladder is your friend",
                      kicker="Part 1 · Foundations", section="Foundations")
add_bullets(slide, [
    (0, [("Scan before the Foley", {"bold": True}),
         (" — a full bladder is the acoustic window", {})]),
    (0, "Transverse and sagittal sweeps; fluid pools posterior to the bladder "
        "(rectovesical space; pouch of Douglas in women)"),
    (0, [("Most dependent peritoneal space in the pelvis", {"bold": True}),
         (" — may be positive when upper quadrants are still negative", {})]),
    (0, [("Pitfalls: ", {"bold": True, "color": RED}),
         ("seminal vesicles, ovarian cysts, small physiologic fluid in "
          "premenopausal women", {})]),
], w=Inches(7.4))
image_placeholder(slide, Inches(8.35), Inches(1.95), Inches(4.35), Inches(2.35),
                  "Pelvic view — transverse")
image_placeholder(slide, Inches(8.35), Inches(4.45), Inches(4.35), Inches(2.35),
                  "Positive pelvic view")
add_notes(slide, "~1.5 min. Quick physiologic-fluid caveat: trace pelvic fluid in a "
                 "young woman is interpreted in clinical context, not reflexively "
                 "called a positive FAST.")

# ---- 11: cardiac
slide = content_slide("Cardiac view — it changes the operation",
                      kicker="Part 1 · Foundations", section="Foundations")
add_bullets(slide, [
    (0, "Subxiphoid four-chamber first; parasternal long-axis as the rescue window"),
    (0, [("Pericardial effusion in trauma = hemopericardium until proven "
          "otherwise", {"bold": True})]),
    (0, "Tamponade physiology: RV diastolic collapse, plethoric IVC"),
    (0, [("In penetrating chest trauma this is the star view", {"bold": True}),
         (" — it moves patients to the OR before they arrest", {})]),
    (0, [("Pitfalls: ", {"bold": True, "color": RED}),
         ("epicardial fat pad; clotted hemopericardium may be isoechoic", {})]),
], w=Inches(7.4))
image_placeholder(slide, Inches(8.35), Inches(1.95), Inches(4.35), Inches(2.35),
                  "Normal subxiphoid clip")
image_placeholder(slide, Inches(8.35), Inches(4.45), Inches(4.35), Inches(2.35),
                  "Pericardial effusion ± RV collapse")
add_notes(slide, "~2 min. If the liver is in the way, push the depth and use it as a "
                 "window; if subxiphoid fails, go parasternal without hesitation.")

# ---- 12: pneumothorax
slide = content_slide("The “e” in eFAST — pneumothorax",
                      kicker="Part 1 · Foundations", section="Foundations")
add_bullets(slide, [
    (0, [("Lung sliding present → no pneumothorax at that interspace", {"bold": True}),
         ("  (B-lines or lung pulse also exclude it there)", {"size": 15, "color": GRAY})]),
    (0, [("Absent sliding alone is NOT diagnostic", {"bold": True, "color": RED}),
         (" — mainstem intubation, blebs, adhesions, apnea all abolish it", {})]),
    (0, [("Lung point = the specific sign", {"bold": True}),
         (" — the edge where sliding reappears", {})]),
    (0, "M-mode memory hook: seashore (normal) vs. barcode/stratosphere (PTX)"),
    (0, [("Supine anterior chest is the money spot", {"bold": True}),
         (" — air rises to the least dependent point", {})]),
], w=Inches(7.4))
image_placeholder(slide, Inches(8.35), Inches(1.95), Inches(4.35), Inches(2.35),
                  "Lung sliding + M-mode seashore")
image_placeholder(slide, Inches(8.35), Inches(4.45), Inches(4.35), Inches(2.35),
                  "Absent sliding / lung point / barcode")
add_notes(slide, "~2.5 min. This is where eFAST clearly beats the supine chest film. "
                 "Live-demo the M-mode signs if a machine is available on stage.")

# ---- 13: hemothorax
slide = content_slide("Hemothorax — look above the diaphragm",
                      kicker="Part 1 · Foundations", section="Foundations")
add_bullets(slide, [
    (0, "Same RUQ/LUQ windows, angled cephalad — no extra probe placement needed"),
    (0, [("Spine sign: ", {"bold": True}),
         ("vertebral bodies visible above the diaphragm = fluid "
          "(aerated lung normally scatters the beam)", {})]),
    (0, "Detects far smaller collections than the supine chest film, on which "
        "blood layers posteriorly and even large hemothoraces can look subtle"),
    (0, "Quantification is possible but rarely changes the trauma decision — "
        "present or absent is what matters tonight"),
], w=Inches(7.4))
image_placeholder(slide, Inches(8.35), Inches(1.95), Inches(4.35), Inches(2.35),
                  "Normal diaphragm — mirror artifact")
image_placeholder(slide, Inches(8.35), Inches(4.45), Inches(4.35), Inches(2.35),
                  "Hemothorax with spine sign")
add_notes(slide, "~1.5 min. The mirror-artifact-vs-spine-sign pairing is intuitive "
                 "even for novices — good audience moment.")

# ---- 14: evidence
slide = content_slide("What the evidence honestly supports",
                      kicker="Part 1 · Foundations", section="Foundations")
rows = [
    ("Free intraperitoneal fluid",
     "Sensitivity varies widely across studies (~60–95%+, operator- and "
     "population-dependent); specificity consistently high (~95%+)", TEAL),
    ("Pneumothorax",
     "US markedly more sensitive than supine CXR in meta-analyses "
     "(pooled sensitivity roughly 85–90% vs ~40–50% for CXR), with comparable "
     "high specificity", TEAL),
    ("Pericardial effusion",
     "High sensitivity and specificity in penetrating trauma series; "
     "false negatives reported with decompression into the pleural space", NAVY),
    ("The honest limitation",
     "A negative eFAST does NOT exclude injury — it excludes the findings it "
     "looks for, at that moment, in those windows", RED),
]
y0 = Inches(1.9)
for i, (head, body, c) in enumerate(rows):
    panel_text(slide, Inches(0.6), y0 + i * Inches(1.22), Inches(12.1), Inches(1.08),
               head, body, accent=c, body_size=14)
add_notes(slide, "~2.5 min. Figures are deliberately given as ranges — exact pooled "
                 "estimates differ by meta-analysis; verify against the sources on "
                 "the reference slide (and your own literature check) before "
                 "presenting. The rhetorical move: concede the limits early so the "
                 "rural/MCI claims later are credible.")

# ---- 15: pitfalls
slide = content_slide("Pitfalls — how eFAST lies to you",
                      kicker="Part 1 · Foundations", section="Foundations")
panel_text(slide, Inches(0.6), Inches(1.95), Inches(5.95), Inches(4.7),
           "False negatives", [
               "Small volumes below detection threshold (~200–500 mL "
               "intraperitoneal, technique-dependent)",
               "Clotted blood — isoechoic, easy to miss",
               "Retroperitoneal hemorrhage — essentially invisible",
               "Hollow viscus and mesenteric injury",
               "Obesity, subcutaneous emphysema, uncooperative patient",
               "Scanning too early — bleeding hasn't accumulated yet",
           ], accent=RED, body_size=14)
panel_text(slide, Inches(6.75), Inches(1.95), Inches(5.95), Inches(3.15),
           "False positives", [
               "Pre-existing ascites (cirrhosis, dialysis, malignancy)",
               "Physiologic pelvic fluid in premenopausal women",
               "Epicardial fat pad mimicking pericardial fluid",
               "Fluid-filled bowel or stomach misread as free fluid",
           ], accent=AMBER, header_color=TEAL_DARK, body_size=14)
pb = panel(slide, Inches(6.75), Inches(5.3), Inches(5.95), Inches(1.35), fill=NAVY)
tb = _box(slide, Inches(7.0), Inches(5.45), Inches(5.5), Inches(1.1))
_text(tb, [("The antidote to every pitfall:\n", {"size": 15, "color": WHITE, "bold": True}),
           ("serial exams + clinical context", {"size": 18, "color": AMBER, "bold": True})],
      anchor=MSO_ANCHOR.MIDDLE)
add_notes(slide, "~2 min. End Part 1 on the serial-exam mantra — it recurs in the "
                 "rural section (repeat during long transports) and the MCI section "
                 "(re-triage loops).")

# ================================================================ PART 2
divider("02", "Trauma decision-making",
        "The same four answers drive different actions — depending on physiology",
        "~30 sec. Part 2 is short and algorithmic: where the exam sits, and how "
        "positive/negative branches differ in stable vs unstable patients.")

# ---- 17: ATLS position
slide = content_slide("Where eFAST sits: inside the primary survey",
                      kicker="Part 2 · Decision-making", section="Decision-making")
add_bullets(slide, [
    (0, [("An adjunct to “C”", {"bold": True}),
         (" in the ATLS® primary survey — performed during resuscitation, "
          "not queued after it", {})]),
    (0, "Answers one question the physical exam cannot: "
        "“where is the blood?” in the unstable blunt-trauma patient"),
    (0, [("Repeatable at zero marginal cost", {"bold": True}),
         (" — repeat after any deterioration, after interventions, "
          "and before transfer", {})]),
    (0, "Documented like any other exam: views obtained, findings, "
        "adequacy, time"),
])
add_notes(slide, "~1.5 min. Position matters: teams that treat eFAST as an early, "
                 "integrated step act on it; teams that treat it as an imaging "
                 "order lose the speed advantage.")

# ---- 18: blunt algorithm
slide = content_slide("Blunt trauma: the 2×2 that runs the room",
                      kicker="Part 2 · Decision-making", section="Decision-making")
cells = [
    ("UNSTABLE + eFAST positive", "Operating room (or transfer to one) — "
     "no CT on the way", RED, WHITE),
    ("UNSTABLE + eFAST negative", "Hunt elsewhere: chest, pelvis, retroperitoneum, "
     "long bones; repeat the exam; consider other shock causes", AMBER, NAVY),
    ("STABLE + eFAST positive", "CT with contrast while stability lasts — "
     "grade the injury, plan selective management", TEAL, WHITE),
    ("STABLE + eFAST negative", "Serial exams ± CT by mechanism and exam — "
     "a single negative never closes the story", NAVY, WHITE),
]
for i, (head, body, fill, hcolor) in enumerate(cells):
    x = Inches(0.6) + (i % 2) * Inches(6.25)
    y = Inches(2.0) + (i // 2) * Inches(2.35)
    p = panel(slide, x, y, Inches(5.95), Inches(2.1), fill=fill)
    tb = _box(slide, x + Inches(0.3), y + Inches(0.2), Inches(5.4), Inches(1.75))
    _text(tb, [(head + "\n", {"size": 16, "bold": True, "color": hcolor}),
               (body, {"size": 14, "color": WHITE if fill in (RED, TEAL, NAVY) else NAVY})])
add_notes(slide, "~2.5 min. Walk the quadrants clockwise from top-left. The "
                 "top-left cell is the entire justification for the exam: it "
                 "removes CT from the path of the dying patient.")

# ---- 19: penetrating
slide = content_slide("Penetrating trauma: different priorities",
                      kicker="Part 2 · Decision-making", section="Decision-making")
add_bullets(slide, [
    (0, [("Pericardial view first", {"bold": True}),
         (" in thoracoabdominal wounds — hemopericardium sends the patient "
          "to the OR immediately", {})]),
    (0, "Positive abdominal windows help choose which cavity to open first "
        "in multi-cavity injury"),
    (0, [("A negative exam means little for hollow viscus, diaphragm, "
          "or trajectory", {"bold": True, "color": RED}),
         (" — penetrating pathways still need exploration, imaging, "
          "or observation per your algorithm", {})]),
    (0, "Serial exams again: delayed hemopericardium and delayed "
        "hemoperitoneum are both described"),
])
add_notes(slide, "~1.5 min. One-line summary for the audience: in penetrating "
                 "trauma, eFAST rules things IN beautifully and rules almost "
                 "nothing OUT.")

# ---- 20: case 1 resolution
slide = content_slide("Back to our motorcyclist", kicker="Case resolution",
                      section="Case")
add_bullets(slide, [
    (0, [("eFAST at 4 minutes after arrival: ", {})]),
    (1, [("anechoic stripe in Morison's pouch + caudal liver tip", {"bold": True})]),
    (1, "pericardium dry, lungs sliding bilaterally"),
    (0, [("Unstable + positive abdomen = the decision is made", {"bold": True, "color": TEAL}),
         (" — activate transfer to the trauma center NOW; "
          "blood en route, no CT attempted", {})]),
    (0, "Receiving surgeon gets a concrete finding, not a hunch — "
        "the OR is ready before the helicopter lands"),
    (0, [("Time from probe-on-skin to transfer decision: under 5 minutes.",
          {"italic": True, "color": GRAY})]),
], y=Inches(1.95), h=Inches(4.6), w=Inches(7.9))
image_placeholder(slide, Inches(8.9), Inches(4.5), Inches(3.8), Inches(2.3),
                  "The actual positive clip, if you have one like it")
add_notes(slide, "~2 min. Land the emotional beat: nothing about this decision "
                 "required a radiologist, a CT, or a big hospital. That is the "
                 "bridge into Part 3.")

# ================================================================ PART 3
divider("03", "The rural reality",
        "Where eFAST is not an adjunct — it is the imaging department",
        "~30 sec. Shift the frame: everything so far assumed a trauma center. "
        "Most of the world's trauma does not happen in one.")

# ---- 22: rural problem
slide = content_slide("The rural trauma problem", kicker="Part 3 · Rural",
                      section="Rural")
probs = [
    ("Distance", "Definitive care is hours away; outcomes worsen with time "
     "to hemorrhage control"),
    ("Imaging gaps", "No CT, or none after hours; plain films and clinical "
     "exam carry the load"),
    ("Staffing", "No in-house surgeon or anesthetist; one clinician may run "
     "the entire resuscitation"),
    ("Transfer stakes", "Under-triage kills; over-triage drains helicopters, "
     "beds, and family finances"),
]
for i, (head, body) in enumerate(probs):
    x = Inches(0.6) + (i % 2) * Inches(6.25)
    y = Inches(2.0) + (i // 2) * Inches(2.2)
    panel_text(slide, x, y, Inches(5.95), Inches(1.95), head, body,
               accent=NAVY if i % 2 else TEAL, body_size=15)
add_notes(slide, "~2 min. Localize this slide: swap in your own region's transfer "
                 "times and CT availability if you have them — audiences engage "
                 "far more with local numbers. [ADD LOCAL DATA IF AVAILABLE]")

# ---- 23: eFAST changes the calculus
slide = content_slide("eFAST changes the rural calculus",
                      kicker="Part 3 · Rural", section="Rural")
add_bullets(slide, [
    (0, [("Triage-to-transfer: ", {"bold": True}),
         ("a positive exam converts “worried” into “documented free fluid” — "
          "and justifies the helicopter", {})]),
    (0, [("Destination selection: ", {"bold": True}),
         ("positive pericardium or abdomen argues for the trauma center, "
          "not the nearest hospital", {})]),
    (0, [("Objective handover: ", {"bold": True}),
         ("findings transmit cleanly by phone — “positive RUQ, dry pericardium” "
          "beats “looks bad”", {})]),
    (0, [("Repeat during transport: ", {"bold": True}),
         ("a negative patient who becomes positive mid-transfer changes "
          "the receiving team's preparation", {})]),
    (0, [("The machine pays for itself in avoided wrong decisions, "
          "not billed studies.", {"italic": True, "color": GRAY})]),
])
add_notes(slide, "~2 min. The handover point is underrated: eFAST creates a shared, "
                 "objective vocabulary between a rural GP and a trauma surgeon "
                 "who have never met.")

# ---- 24: handheld & tele
slide = content_slide("Handhelds and tele-ultrasound",
                      kicker="Part 3 · Rural", section="Rural")
add_bullets(slide, [
    (0, [("Pocket-sized probes", {"bold": True}),
         (" now deliver eFAST-adequate images at a fraction of cart prices — "
          "battery-powered, phone-connected", {})]),
    (0, [("Tele-ultrasound: ", {"bold": True}),
         ("live image streaming + remote expert guidance of a novice's hand "
          "is feasible over ordinary mobile networks", {})]),
    (0, "Prehospital eFAST in ambulances and HEMS is reported and feasible; "
        "the challenge is training and governance, not physics"),
    (0, [("Caveats: ", {"bold": True, "color": RED}),
         ("image quality varies by device; battery and disinfection logistics "
          "are real; a handheld does not confer skill", {})]),
], w=Inches(7.6))
image_placeholder(slide, Inches(8.5), Inches(1.95), Inches(4.2), Inches(2.3),
                  "Handheld probe photo (your device)")
image_placeholder(slide, Inches(8.5), Inches(4.4), Inches(4.2), Inches(2.3),
                  "Tele-ultrasound screenshot / setup")
add_notes(slide, "~2 min. If you have personal experience with a specific handheld "
                 "or tele-US program, tell that story here — first-person beats "
                 "citation. Avoid naming brands unless disclosed.")

# ---- 25: rural workflow
slide = content_slide("A rural workflow that survives real life",
                      kicker="Part 3 · Rural", section="Rural")
steps = [
    ("On arrival", "eFAST inside the primary survey — before the Foley, "
     "before the log-roll if feasible"),
    ("Decision point", "Unstable + positive → transfer/OR pathway immediately; "
     "unstable + negative → widen the search, repeat in 10–15 min"),
    ("Pre-departure", "Repeat exam; document views, findings, time — "
     "send images with the patient if you can"),
    ("In transit", "Repeat after deterioration or long legs of transport; "
     "relay changes ahead"),
]
y0 = Inches(1.95)
for i, (head, body) in enumerate(steps):
    y = y0 + i * Inches(1.22)
    tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y,
                                 Inches(2.5), Inches(1.05))
    _set_fill(tag, TEAL if i % 2 == 0 else NAVY)
    _text(tag, head, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
          anchor=MSO_ANCHOR.MIDDLE)
    bb = _box(slide, Inches(3.4), y + Inches(0.05), Inches(9.3), Inches(1.0))
    _text(bb, body, size=15, color=SLATE, anchor=MSO_ANCHOR.MIDDLE)
add_notes(slide, "~2 min. Emphasize documentation: in transfer systems, the exam "
                 "that isn't written down (or sent) didn't happen.")

# ---- 26: rural case
slide = content_slide("Case 2 — the tractor rollover", kicker="Part 3 · Rural",
                      section="Case")
add_bullets(slide, [
    (0, "58-year-old farmer, tractor rollover, 2 hours from the trauma center"),
    (1, "Initially stable: BP 118/76, HR 96 — “he looks okay”"),
    (0, [("First eFAST: negative.", {"bold": True}),
         ("  Kept for observation with hourly vitals", {})]),
    (0, [("Repeat at 45 min (new tachycardia): ", {})]),
    (1, [("free fluid now visible at the caudal liver tip", {"bold": True, "color": RED})]),
    (0, "Transfer launched before hypotension ever appeared — splenic injury "
        "managed operatively that evening"),
    (0, [("Moral: the serial exam is the rural safety net.",
          {"bold": True, "color": TEAL})]),
], y=Inches(1.9))
add_notes(slide, "~2 min. This case sells serial exams better than any test "
                 "characteristic can. Swap in a real (de-identified) case of your "
                 "own if available — it will land harder.")

# ---- 27: training
slide = content_slide("Keeping skills alive in low-volume settings",
                      kicker="Part 3 · Rural", section="Rural")
add_bullets(slide, [
    (0, "Learning curve to competent basic eFAST is modest — commonly cited "
        "thresholds are in the tens of proctored exams, not hundreds"),
    (0, [("The rural problem is decay, not acquisition: ", {"bold": True}),
         ("low trauma volume erodes image-generation skill", {})]),
    (0, [("Countermeasures that work in practice:", {"bold": True})]),
    (1, "scan non-trauma patients deliberately (normal anatomy reps)"),
    (1, "periodic image-review sessions with a remote expert (QA loop)"),
    (1, "simulation and phantom refreshers; tele-mentored live exams"),
    (0, "Credentialing: follow your national/institutional POCUS framework — "
        "and record your numbers"),
])
add_notes(slide, "~2 min. Practical, non-preachy tone. If your hospital runs a QA "
                 "or tele-review loop, describe it concretely. [ADD LOCAL "
                 "PROGRAM DETAILS IF APPLICABLE]")

# ================================================================ PART 4
divider("04", "Mass-casualty incidents",
        "When needs exceed resources, the probe becomes a triage instrument",
        "~30 sec. Reframe once more: same exam, but now the scarce resource is "
        "not CT — it is time, operating rooms, and transport slots.")

# ---- 29: MCI problem
slide = content_slide("The MCI triage problem", kicker="Part 4 · Mass casualty",
                      section="MCI")
add_bullets(slide, [
    (0, [("Definition of the problem: ", {"bold": True}),
         ("casualties exceed the resources available to treat them "
          "at the usual standard", {})]),
    (0, "First-pass triage (START and relatives) is physiologic — it sorts by "
        "vitals and mobility, and is blind to internal hemorrhage"),
    (0, [("The dangerous group: ", {"bold": True, "color": RED}),
         ("“delayed” patients who are quietly bleeding — physiology "
          "compensates until it doesn't", {})]),
    (0, "CT is a single-lane bridge in a surge; whatever sorts patients "
        "without it protects the whole system"),
])
add_notes(slide, "~2 min. Set up the gap precisely: MCI triage tools measure "
                 "compensation, not injury. eFAST looks one layer deeper for "
                 "four specific, actionable findings.")

# ---- 30: MCI evidence
slide = content_slide("What disasters have taught us",
                      kicker="Part 4 · Mass casualty", section="MCI")
rows = [
    ("1988 Armenia earthquake",
     "Sonographic screening of hundreds of casualties for abdominal injury — "
     "an early landmark of mass ultrasound triage (Sarkisian et al.)", TEAL),
    ("Subsequent earthquakes & conflicts",
     "Repeated reports (e.g., Marmara 1999, Wenchuan 2008, Haiti 2010; military "
     "deployments) of FAST/eFAST used for surge triage where CT was absent "
     "or overwhelmed", NAVY),
    ("Consistent themes across reports",
     "Fast per-patient exam times, usable in austere conditions, helps "
     "prioritize surgery and evacuation; performance depends on operator "
     "skill and re-examination", TEAL),
]
y0 = Inches(1.95)
for i, (head, body, c) in enumerate(rows):
    panel_text(slide, Inches(0.6), y0 + i * Inches(1.55), Inches(12.1), Inches(1.4),
               head, body, accent=c, body_size=14)
add_notes(slide, "~2.5 min. IMPORTANT: verify the disaster-literature details "
                 "(authors, casualty numbers, exact findings) against the sources "
                 "before presenting — the slide deliberately avoids specific "
                 "performance numbers. The robust, defensible claim: ultrasound "
                 "triage in disasters is repeatedly reported as feasible and "
                 "operationally useful, not that it has trial-grade evidence.")

# ---- 31: MCI integration
slide = content_slide("Where eFAST plugs into MCI triage",
                      kicker="Part 4 · Mass casualty", section="MCI")
steps = [
    ("Primary triage", "START (or local equivalent) sorts by physiology — "
     "seconds per patient, no equipment", NAVY),
    ("Secondary triage + eFAST", "Sonographic screen of “delayed” and "
     "borderline “immediate” patients — find the quiet bleeders", TEAL),
    ("Resource assignment", "Positive abdomen/pericardium → surgical queue or "
     "first evacuation slots; negative → serial re-triage loop", AMBER),
]
for i, (head, body, c) in enumerate(steps):
    x = Inches(0.6) + i * Inches(4.28)
    p = panel(slide, x, Inches(2.1), Inches(3.95), Inches(3.3),
              fill=LIGHT, accent=c)
    tb = _box(slide, x + Inches(0.25), Inches(2.3), Inches(3.5), Inches(2.9))
    _text(tb, [(f"{i+1}\n", {"size": 34, "bold": True, "color": c}),
               (head + "\n", {"size": 17, "bold": True, "color": NAVY}),
               (body, {"size": 13.5, "color": SLATE})])
    if i < 2:
        ar = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(3.97),
                                    Inches(3.5), Inches(0.31), Inches(0.5))
        _set_fill(ar, GRAY)
tb = _box(slide, Inches(0.62), Inches(5.7), Inches(12.0), Inches(1.0))
_text(tb, [("Key principle: ", {"size": 17, "bold": True, "color": TEAL}),
           ("eFAST refines triage categories — it does not replace the triage "
            "system or delay obvious salvage.", {"size": 17, "color": SLATE})])
add_notes(slide, "~2 min. Guard against the failure mode where ultrasound becomes "
                 "a bottleneck: it is a screen inside secondary triage, never a "
                 "gate in front of primary triage.")

# ---- 32: MCI workflow
slide = content_slide("A practical MCI ultrasound workflow",
                      kicker="Part 4 · Mass casualty", section="MCI")
add_bullets(slide, [
    (0, [("One trained operator roams with one machine", {"bold": True}),
         (" — the sonographer goes to the patient, never the reverse", {})]),
    (0, "Abbreviated exam, ~1–2 minutes per patient: abdomen + pericardium "
        "first; thorax when it changes disposition"),
    (0, [("Mark the result on the patient", {"bold": True}),
         (" — triage tag or skin marker (e.g., “F+ RUQ 14:32”), because charts "
          "do not exist in an MCI", {})]),
    (0, "Re-triage loop: negatives get rescanned on a schedule or on any "
        "physiologic change"),
    (0, "Rehearse it: include ultrasound in your hospital's next MCI drill — "
        "roles, batteries, disinfection, tag conventions"),
])
add_notes(slide, "~2.5 min. The skin-marking detail always draws nods — it converts "
                 "the abstract 'ultrasound triage' into something an audience can "
                 "operationalize on Monday.")

# ---- 33: MCI limitations
slide = content_slide("Limitations under surge conditions — be honest",
                      kicker="Part 4 · Mass casualty", section="MCI")
add_bullets(slide, [
    (0, [("A single negative scan is weakest exactly when repeat scans are "
          "hardest to deliver", {"bold": True, "color": RED})]),
    (0, "Operator fatigue and cognitive load degrade accuracy over hours "
        "of continuous scanning"),
    (0, "One machine = one queue; device failure or dead batteries can "
        "silently remove the capability"),
    (0, "Environmental reality: light, noise, weather, contaminated fields, "
        "probe disinfection between patients"),
    (0, [("Mitigations: ", {"bold": True, "color": TEAL}),
         ("scheduled operator rotation, spare batteries and probes, "
          "pre-agreed abbreviated protocol, drilled tag conventions", {})]),
])
add_notes(slide, "~1.5 min. Credibility slide — conceding the failure modes makes "
                 "the overall MCI argument stronger, not weaker.")

# ---- 34: combined case
slide = content_slide("Case 3 — the bus rollover, 30 km out of town",
                      kicker="Part 4 · Mass casualty", section="Case")
add_bullets(slide, [
    (0, "Interurban bus rollover: 23 casualties arrive at a district hospital "
        "over 40 minutes; one general surgeon, one OR, no CT after hours"),
    (0, "START triage: 4 immediate, 9 delayed, 8 minor, 2 deceased"),
    (0, [("eFAST on the 4 immediate + 9 delayed (≈20 minutes total):", {"bold": True})]),
    (1, [("2 positive abdomens found among the “delayed”", {"bold": True, "color": RED}),
         (" — both upgraded, one to the local OR, one to the first helicopter", {})]),
    (1, "1 “immediate” with dry belly and pericardium → hemothorax found → "
        "chest drain locally, evacuation deferred safely"),
    (0, [("Same probe, three roles: rural imaging, triage refinement, "
          "resource allocation.", {"bold": True, "color": TEAL})]),
], y=Inches(1.9))
add_notes(slide, "~2.5 min. Composite teaching case — present it as such. It "
                 "deliberately fuses Parts 3 and 4: limited hospital + surge. "
                 "Walk it slowly; this is the talk's synthesis moment.")

# ================================================================ PART 5
divider("05", "The road ahead",
        "Training pipelines, tele-mentoring, AI assistance — and five closing points",
        "~20 sec. Final stretch: keep energy up, this section is short.")

# ---- 36: future
slide = content_slide("Where this is heading", kicker="Part 5 · Future",
                      section="Future")
add_bullets(slide, [
    (0, [("AI-assisted interpretation: ", {"bold": True}),
         ("automated free-fluid and lung-sliding detection is in active "
          "development; early tools aim to support novices, not replace "
          "training", {})]),
    (0, [("Tele-mentoring at scale: ", {"bold": True}),
         ("remote guidance can extend one expert across many rural sites "
          "and disaster scenes", {})]),
    (0, [("Cheaper, tougher hardware: ", {"bold": True}),
         ("handhelds are approaching commodity pricing — the barrier is "
          "shifting from equipment to education", {})]),
    (0, [("Systems integration: ", {"bold": True}),
         ("images and findings flowing into transfer networks and disaster "
          "command dashboards", {})]),
    (0, [("Treat AI claims with the same skepticism as any diagnostic claim — "
          "ask for the validation data.", {"italic": True, "color": GRAY})]),
])
add_notes(slide, "~2 min. Stay grounded — this audience will include skeptics. "
                 "The defensible line: the bottleneck is moving from hardware "
                 "to training and governance, and AI may compress the training "
                 "curve if validated properly.")

# ---- 37: take-homes
slide = content_slide("Five things to take home", kicker="Closing",
                      section="Closing")
points = [
    "eFAST answers four binary questions in minutes — it is a disposition "
    "tool, not a survey of injury",
    "Unstable + positive abdomen = operating room; no CT belongs in "
    "that pathway",
    "A negative exam never closes the story — the serial exam is the "
    "safety net, especially before and during rural transfer",
    "In mass casualty, eFAST refines secondary triage: it finds the quiet "
    "bleeders that physiologic triage misses",
    "The barrier is no longer the machine — it is training, maintenance "
    "of skill, and rehearsed workflows",
]
y0 = Inches(1.9)
for i, txt in enumerate(points):
    y = y0 + i * Inches(0.98)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.65), y + Inches(0.1),
                                 Inches(0.62), Inches(0.62))
    _set_fill(dot, TEAL if i % 2 == 0 else NAVY)
    nb = _box(slide, Inches(0.65), y + Inches(0.13), Inches(0.62), Inches(0.55))
    _text(nb, str(i + 1), size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    tb = _box(slide, Inches(1.55), y + Inches(0.02), Inches(11.2), Inches(0.9))
    _text(tb, txt, size=16.5, color=SLATE, anchor=MSO_ANCHOR.MIDDLE)
add_notes(slide, "~2 min. Read all five slowly — this is the slide people "
                 "photograph. Then transition to thanks/questions.")

# ---- 38: references
slide = content_slide("Selected reading", kicker="References", section="Closing")
add_bullets(slide, [
    (0, "Scalea TM, et al. Focused Assessment with Sonography for Trauma (FAST): "
        "results of an international consensus conference. J Trauma. 1999."),
    (0, "Rozycki GS, et al. Surgeon-performed ultrasound for the assessment of "
        "truncal injuries. Ann Surg. 1998."),
    (0, "Kirkpatrick AW, et al. Hand-held thoracic sonography for detecting "
        "post-traumatic pneumothoraces: the extended FAST (EFAST). J Trauma. 2004."),
    (0, "Sarkisian AE, et al. Sonographic screening of mass casualties for "
        "abdominal and renal injuries following the 1988 Armenian earthquake. "
        "J Trauma. 1991."),
    (0, "American College of Surgeons. Advanced Trauma Life Support (ATLS), "
        "10th edition. 2018."),
    (0, [("[VERIFY all citation details against the primary sources before "
          "presenting; add recent meta-analyses of your choosing]",
          {"italic": True, "color": RED, "size": 14})]),
], size=15, gap=12)
add_notes(slide, "Do not read this slide aloud. Before the talk: verify every "
                 "citation (authors, journal, year) against PubMed and add 1–2 "
                 "recent meta-analyses on eFAST test performance that you have "
                 "personally checked.")

# ---- 39: thank you
slide = prs.slides.add_slide(BLANK)
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
_set_fill(bg, NAVY)
stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.9), SLIDE_W, Inches(0.12))
_set_fill(stripe, TEAL)
tb = _box(slide, Inches(0.95), Inches(2.3), Inches(11.5), Inches(1.3))
_text(tb, "Thank you.", size=60, color=WHITE, bold=True)
sb = _box(slide, Inches(1.0), Inches(3.8), Inches(11.2), Inches(0.8))
_text(sb, "Questions — and your own rural or MCI experiences — very welcome.",
      size=22, color=RGBColor(0xB8, 0xC9, 0xD9))
cb = _box(slide, Inches(1.0), Inches(4.9), Inches(11), Inches(0.8))
_text(cb, [("Chia-Ching Chen, MD   ·   ", {"size": 16, "color": WHITE, "bold": True}),
           ("[EMAIL / CONTACT]", {"size": 16, "color": TEAL, "italic": True})])
add_footer(slide)
add_notes(slide, "Invite audience war stories explicitly — rural and disaster "
                 "ultrasound talks reliably generate good Q&A when you ask for "
                 "experiences, not just questions. Budget ~10 minutes.")

# ---------------------------------------------------------------- save
out = __file__.rsplit("/", 1)[0] + "/eFAST_trauma_keynote.pptx"
prs.save(out)
print(f"Saved {out} with {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
